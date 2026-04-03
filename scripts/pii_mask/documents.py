"""Office document masker — PII replacement in DOCX, XLSX, PPTX.

Uses python-docx, openpyxl, python-pptx for high-level text access
(preserves formatting), then a secondary ZIP+XML pass for content
the high-level APIs miss (footnotes, endnotes, SmartArt, app properties).

Pattern mirrors eml.py: public mask_*() functions take raw bytes +
scanner, return masked bytes.  Called by maskers via BaseMasker's
_mask_document_file() helper.
"""

import io
import logging
import re
import uuid
import zipfile

from lxml import etree

from scripts.pii_mask.scanner import TextScanner

log = logging.getLogger(__name__)

# Office doc extensions we handle
OFFICE_EXTENSIONS = frozenset({".docx", ".xlsx", ".pptx"})

# Separator for batched scanning — random hex string guaranteed not to
# appear in real text or be detected as PII.  Generated once at import
# time, reused for all files in the process.
_BATCH_SEP = f"\n\n{uuid.uuid4().hex}\n\n"

# XML namespaces
_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DC_NS = "http://purl.org/dc/elements/1.1/"
_CP_NS = "http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
_EP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"


# -- Public API ------------------------------------------------------------ #

def mask_docx(raw_bytes: bytes, scanner: TextScanner) -> bytes:
    """Mask PII in a DOCX file.  Returns masked bytes.

    Collects all paragraph texts, scans them in a single batched NER
    call, then distributes masked text back to runs.
    """
    from docx import Document

    doc = Document(io.BytesIO(raw_bytes))

    # Collect all paragraphs from body, headers/footers, comments
    paras = []
    for item in doc.iter_inner_content():
        if hasattr(item, "runs"):
            paras.append(item)
        else:
            _collect_table_paras(item, paras)

    for section in doc.sections:
        for hf in (section.header, section.footer,
                    section.even_page_header, section.even_page_footer,
                    section.first_page_header, section.first_page_footer):
            if hf.is_linked_to_previous:
                continue
            for item in hf.iter_inner_content():
                if hasattr(item, "runs"):
                    paras.append(item)
                else:
                    _collect_table_paras(item, paras)

    comment_authors = []
    if hasattr(doc, "comments"):
        for comment in doc.comments:
            if hasattr(comment, "author") and comment.author:
                comment_authors.append(comment)
            for para in comment.paragraphs:
                paras.append(para)

    # Build list of (joined_text, runs) for paragraphs with content
    segments = []
    for para in paras:
        runs = para.runs
        if not runs:
            continue
        texts = [r.text or "" for r in runs]
        joined = "".join(texts)
        if len(joined) >= 3:
            segments.append((joined, runs, texts))

    # Batch scan all paragraph texts in one NER call
    if segments:
        originals = [s[0] for s in segments]
        masked_texts = _batch_scan(originals, scanner)
        for (orig, runs, run_texts), masked in zip(segments, masked_texts):
            if masked != orig:
                _distribute_to_runs(runs, run_texts, orig, masked)

    # Comment authors — small count, scan individually
    for comment in comment_authors:
        comment.author = scanner.scan(comment.author)

    # Core properties
    _mask_core_properties(doc.core_properties, scanner)

    buf = io.BytesIO()
    doc.save(buf)
    primary_bytes = buf.getvalue()

    return _secondary_xml_pass(primary_bytes, scanner, fmt="docx")


def mask_xlsx(raw_bytes: bytes, scanner: TextScanner) -> bytes:
    """Mask PII in an XLSX file.  Returns masked bytes.

    Collects all string cell values, scans them in a single batched NER
    call, then writes masked values back.  Reduces NER calls from
    O(cells) to O(1) — typically 100-500x faster for large spreadsheets.
    """
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw_bytes))

    # Collect all string cells + comments across all sheets
    cell_entries = []    # (cell, original_value)
    comment_entries = [] # (comment, "text"|"author", original_value)

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and len(cell.value) >= 3:
                    cell_entries.append((cell, cell.value))
                if cell.comment:
                    if cell.comment.text and len(cell.comment.text) >= 3:
                        comment_entries.append(
                            (cell.comment, "text", cell.comment.text))
                    if cell.comment.author and len(cell.comment.author) >= 3:
                        comment_entries.append(
                            (cell.comment, "author", cell.comment.author))

    # Batch scan all values in one NER call
    all_originals = ([v for _, v in cell_entries]
                     + [v for _, _, v in comment_entries])

    if all_originals:
        all_masked = _batch_scan(all_originals, scanner)

        # Apply cell values
        n_cells = len(cell_entries)
        for i, (cell, orig) in enumerate(cell_entries):
            cell.value = all_masked[i]

        # Apply comment text/author
        for i, (comment, attr, orig) in enumerate(comment_entries):
            setattr(comment, attr, all_masked[n_cells + i])

    # Sheet titles — small count, scan individually
    for ws in wb.worksheets:
        if ws.title and len(ws.title) >= 3:
            ws.title = scanner.scan(ws.title)

    # Core properties
    _mask_wb_properties(wb.properties, scanner)

    buf = io.BytesIO()
    wb.save(buf)
    primary_bytes = buf.getvalue()

    return _secondary_xml_pass(primary_bytes, scanner, fmt="xlsx")


def mask_pptx(raw_bytes: bytes, scanner: TextScanner) -> bytes:
    """Mask PII in a PPTX file.  Returns masked bytes.

    Collects all paragraph texts from shapes, notes, and tables,
    scans them in a single batched NER call, then distributes back.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(raw_bytes))

    # Collect all paragraphs from all shapes and notes
    paras = []
    for slide in prs.slides:
        for shape in slide.shapes:
            _collect_shape_paras(shape, paras, MSO_SHAPE_TYPE)
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for para in notes_tf.paragraphs:
                    paras.append(para)

    # Build segments
    segments = []
    for para in paras:
        runs = para.runs
        if not runs:
            continue
        texts = [r.text or "" for r in runs]
        joined = "".join(texts)
        if len(joined) >= 3:
            segments.append((joined, runs, texts))

    # Batch scan
    if segments:
        originals = [s[0] for s in segments]
        masked_texts = _batch_scan(originals, scanner)
        for (orig, runs, run_texts), masked in zip(segments, masked_texts):
            if masked != orig:
                _distribute_to_runs(runs, run_texts, orig, masked)

    # Core properties
    _mask_core_properties(prs.core_properties, scanner)

    buf = io.BytesIO()
    prs.save(buf)
    primary_bytes = buf.getvalue()

    return _secondary_xml_pass(primary_bytes, scanner, fmt="pptx")


def is_office_doc(key: str) -> bool:
    """Check whether an S3 key is a supported Office document."""
    lower = key.lower()
    return (lower.endswith(".docx") or lower.endswith(".xlsx")
            or lower.endswith(".pptx"))


# -- Batched scanning ------------------------------------------------------ #

def _batch_scan(texts: list[str], scanner: TextScanner) -> list[str]:
    """Scan multiple texts in a single NER call.

    Concatenates all texts with a unique separator, runs scanner.scan()
    once, then splits back.  Falls back to per-text scanning if the
    split produces an unexpected number of segments (safety net).

    For N texts this reduces NER calls from N to 1 — the dominant
    cost for spreadsheets with hundreds of cells.
    """
    if not texts:
        return []
    if len(texts) == 1:
        return [scanner.scan(texts[0])]

    joined = _BATCH_SEP.join(texts)
    masked_joined = scanner.scan(joined)
    masked_parts = masked_joined.split(_BATCH_SEP)

    if len(masked_parts) == len(texts):
        return masked_parts

    # Fallback: separator was corrupted by replacement (very unlikely).
    # Scan each text individually.
    log.warning("Batch scan split mismatch (%d vs %d), falling back to "
                "per-text scan", len(masked_parts), len(texts))
    return [scanner.scan(t) for t in texts]


# -- Internal helpers ------------------------------------------------------ #

def _distribute_to_runs(runs, run_texts: list[str],
                        original: str, masked: str) -> None:
    """Distribute masked text back across runs preserving formatting."""
    if len(masked) == len(original):
        pos = 0
        for run, orig_text in zip(runs, run_texts):
            run.text = masked[pos:pos + len(orig_text)]
            pos += len(orig_text)
    else:
        runs[0].text = masked
        for run in runs[1:]:
            run.text = ""


def _collect_table_paras(table, paras: list) -> None:
    """Recursively collect all paragraphs from a table."""
    for row in table.rows:
        for cell in row.cells:
            for item in cell.iter_inner_content():
                if hasattr(item, "runs"):
                    paras.append(item)
                else:
                    _collect_table_paras(item, paras)


def _collect_shape_paras(shape, paras: list, MSO_SHAPE_TYPE) -> None:
    """Recursively collect all paragraphs from a PPTX shape."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_shape_paras(child, paras, MSO_SHAPE_TYPE)
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            paras.append(para)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    paras.append(para)


def _mask_core_properties(cp, scanner: TextScanner) -> None:
    """Mask core document properties (author, title, etc.)."""
    for attr in ("author", "last_modified_by", "subject", "title",
                 "keywords", "comments", "description"):
        val = getattr(cp, attr, None)
        if isinstance(val, str) and len(val) >= 3:
            setattr(cp, attr, scanner.scan(val))


def _mask_wb_properties(props, scanner: TextScanner) -> None:
    """Mask openpyxl workbook properties (different attribute names)."""
    for attr in ("creator", "lastModifiedBy", "subject", "title",
                 "keywords", "description"):
        val = getattr(props, attr, None)
        if isinstance(val, str) and len(val) >= 3:
            setattr(props, attr, scanner.scan(val))


# -- Secondary XML pass ---------------------------------------------------- #

# XML paths to scan per format.
# Each entry: (path_or_glob, namespace_uri, element_local_name)
_DOCX_XML_TARGETS = [
    ("word/footnotes.xml", _W_NS, "t"),
    ("word/endnotes.xml", _W_NS, "t"),
]
_PPTX_XML_TARGETS = [
    ("ppt/diagrams/data*.xml", _A_NS, "t"),
]
_APP_PROPS_PATH = "docProps/app.xml"


def _secondary_xml_pass(raw_bytes: bytes, scanner: TextScanner,
                        fmt: str) -> bytes:
    """Scan XML files inside the ZIP that high-level libraries miss."""
    targets = []
    if fmt == "docx":
        targets = _DOCX_XML_TARGETS
    elif fmt == "pptx":
        targets = _PPTX_XML_TARGETS

    in_buf = io.BytesIO(raw_bytes)
    out_buf = io.BytesIO()

    with zipfile.ZipFile(in_buf, "r") as zin, \
         zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            # Check against targets (supports glob patterns)
            for path_pattern, ns, elem_name in targets:
                if _match_path(item.filename, path_pattern):
                    data = _mask_xml_text_elements(
                        data, ns, elem_name, scanner)
                    break

            # App properties (all formats)
            if item.filename == _APP_PROPS_PATH:
                data = _mask_app_properties(data, scanner)

            zout.writestr(item, data)

    return out_buf.getvalue()


def _match_path(filename: str, pattern: str) -> bool:
    """Match a ZIP entry path against a pattern (supports * glob)."""
    if "*" not in pattern:
        return filename == pattern
    regex = re.escape(pattern).replace(r"\*", ".*")
    return re.fullmatch(regex, filename) is not None


def _mask_xml_text_elements(xml_bytes: bytes, ns: str,
                            elem_name: str,
                            scanner: TextScanner) -> bytes:
    """Replace text in all matching XML elements."""
    try:
        root = etree.fromstring(xml_bytes)
    except etree.XMLSyntaxError:
        return xml_bytes  # corrupted XML — pass through

    for elem in root.iter(f"{{{ns}}}{elem_name}"):
        if elem.text and len(elem.text) >= 3:
            elem.text = scanner.scan(elem.text)

    return etree.tostring(root, xml_declaration=True,
                          encoding="UTF-8", standalone=True)


def _mask_app_properties(xml_bytes: bytes,
                         scanner: TextScanner) -> bytes:
    """Mask company/manager in docProps/app.xml."""
    try:
        root = etree.fromstring(xml_bytes)
    except etree.XMLSyntaxError:
        return xml_bytes

    # App properties use the extended-properties namespace.
    # Tags: Company, Manager, HyperlinkBase
    for tag in ("Company", "Manager", "HyperlinkBase"):
        for elem in root.iter(f"{{{_EP_NS}}}{tag}"):
            if elem.text and len(elem.text) >= 3:
                elem.text = scanner.scan(elem.text)
        # Some files use bare tags without namespace
        for elem in root.iter(tag):
            if elem.text and len(elem.text) >= 3:
                elem.text = scanner.scan(elem.text)

    return etree.tostring(root, xml_declaration=True,
                          encoding="UTF-8", standalone=True)
