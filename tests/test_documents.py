"""Tests for scripts.pii_mask.documents — Office document masking."""

import io

import pytest

from scripts.pii_mask.pii_store import PIIStore
from scripts.pii_mask.scanner import TextScanner


# -- Fixtures -------------------------------------------------------------- #

@pytest.fixture
def store(tmp_path):
    s = PIIStore(str(tmp_path / "test.db"))
    s.add_domain("org_name.com", "example.com")
    s.get_or_create("EMAIL_ADDRESS", "john.doe@org_name.com")
    s.get_or_create("PERSON", "John Doe")
    return s


@pytest.fixture(scope="module")
def _analyzer():
    from presidio_analyzer import AnalyzerEngine
    return AnalyzerEngine()


@pytest.fixture
def scanner(store, _analyzer):
    s = TextScanner.__new__(TextScanner)
    s._store = store
    s._threshold = 0.5
    s._analyzer = _analyzer
    return s


# -- Helpers --------------------------------------------------------------- #

def _make_docx(paragraphs: list[str], author: str = "") -> bytes:
    """Create a minimal DOCX with given paragraph texts."""
    from docx import Document
    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    if author:
        doc.core_properties.author = author
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_docx_with_table(rows: list[list[str]]) -> bytes:
    """Create a DOCX with a table."""
    from docx import Document
    doc = Document()
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            table.rows[i].cells[j].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(cells: dict[str, str], creator: str = "") -> bytes:
    """Create minimal XLSX. cells maps 'A1' -> value."""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for coord, val in cells.items():
        ws[coord] = val
    if creator:
        wb.properties.creator = creator
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_xlsx_with_comment(cell_value: str, comment_text: str,
                            comment_author: str) -> bytes:
    """Create XLSX with a cell comment."""
    from openpyxl import Workbook
    from openpyxl.comments import Comment
    wb = Workbook()
    ws = wb.active
    ws["A1"] = cell_value
    ws["A1"].comment = Comment(comment_text, comment_author)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slide_texts: list[str], author: str = "") -> bytes:
    """Create minimal PPTX with one text box per slide."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1),
                                          Inches(5), Inches(1))
        txBox.text_frame.text = text
    if author:
        prs.core_properties.author = author
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_notes(slide_text: str, notes_text: str) -> bytes:
    """Create PPTX with speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1),
                                      Inches(5), Inches(1))
    txBox.text_frame.text = slide_text
    # Add notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_table(rows: list[list[str]]) -> bytes:
    """Create PPTX with a table on one slide."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        len(rows), len(rows[0]),
        Inches(1), Inches(1), Inches(6), Inches(3))
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            table_shape.table.rows[i].cells[j].text = cell_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# -- DOCX tests ----------------------------------------------------------- #

class TestMaskDocx:
    def test_name_in_paragraph_masked(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx(["Meeting with John Doe tomorrow"])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        text = doc.paragraphs[0].text
        assert "John Doe" not in text
        assert "Meeting with" in text

    def test_email_in_paragraph_masked(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx(["Contact john.doe@org_name.com for details"])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        text = doc.paragraphs[0].text
        assert "john.doe@org_name.com" not in text
        assert "org_name.com" not in text

    def test_no_pii_preserved(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx(["Quarterly revenue report Q4 2025"])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        assert doc.paragraphs[0].text == "Quarterly revenue report Q4 2025"

    def test_core_properties_masked(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx(["Content"], author="John Doe")
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        assert "John Doe" not in (doc.core_properties.author or "")

    def test_table_content_masked(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx_with_table([
            ["Department", "Contact"],
            ["John Doe", "john.doe@org_name.com"],
        ])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        table = doc.tables[0]
        assert "John Doe" not in table.rows[1].cells[0].text
        assert "john.doe@org_name.com" not in table.rows[1].cells[1].text
        # Headers preserved
        assert table.rows[0].cells[0].text == "Department"

    def test_multiple_paragraphs(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx([
            "First paragraph without PII",
            "Second paragraph mentions John Doe",
            "Third paragraph is clean too",
        ])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        assert doc.paragraphs[0].text == "First paragraph without PII"
        assert "John Doe" not in doc.paragraphs[1].text
        assert doc.paragraphs[2].text == "Third paragraph is clean too"

    def test_empty_docx(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        raw = _make_docx([])
        masked = mask_docx(raw, scanner)
        from docx import Document
        doc = Document(io.BytesIO(masked))
        # Should not crash — empty doc has no paragraphs to process
        assert len(doc.paragraphs) == 0


# -- XLSX tests ----------------------------------------------------------- #

class TestMaskXlsx:
    def test_name_in_cell_masked(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        raw = _make_xlsx({"A1": "John Doe", "B1": "john.doe@org_name.com"})
        masked = mask_xlsx(raw, scanner)
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(masked))
        ws = wb.active
        assert "John Doe" not in str(ws["A1"].value)
        assert "john.doe@org_name.com" not in str(ws["B1"].value)

    def test_no_pii_preserved(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        raw = _make_xlsx({"A1": "Revenue", "B1": "Summary Report"})
        masked = mask_xlsx(raw, scanner)
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(masked))
        ws = wb.active
        assert ws["A1"].value == "Revenue"
        assert ws["B1"].value == "Summary Report"

    def test_numeric_cells_unchanged(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws["A1"] = 42
        ws["B1"] = 3.14
        buf = io.BytesIO()
        wb.save(buf)

        masked = mask_xlsx(buf.getvalue(), scanner)
        from openpyxl import load_workbook
        wb2 = load_workbook(io.BytesIO(masked))
        ws2 = wb2.active
        assert ws2["A1"].value == 42
        assert ws2["B1"].value == 3.14

    def test_core_properties_masked(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        raw = _make_xlsx({"A1": "data"}, creator="John Doe")
        masked = mask_xlsx(raw, scanner)
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(masked))
        assert "John Doe" not in (wb.properties.creator or "")

    def test_cell_comment_masked(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        raw = _make_xlsx_with_comment(
            cell_value="data",
            comment_text="Review by John Doe",
            comment_author="John Doe",
        )
        masked = mask_xlsx(raw, scanner)
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(masked))
        ws = wb.active
        comment = ws["A1"].comment
        assert comment is not None
        assert "John Doe" not in comment.text
        assert "John Doe" not in comment.author

    def test_short_strings_skipped(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        raw = _make_xlsx({"A1": "OK", "B1": "No"})
        masked = mask_xlsx(raw, scanner)
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(masked))
        ws = wb.active
        # Strings < 3 chars should pass through unchanged
        assert ws["A1"].value == "OK"
        assert ws["B1"].value == "No"

    def test_multiple_sheets(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        from openpyxl import Workbook
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Contacts"
        ws1["A1"] = "John Doe"
        ws2 = wb.create_sheet("Data")
        ws2["A1"] = "john.doe@org_name.com"
        buf = io.BytesIO()
        wb.save(buf)

        masked = mask_xlsx(buf.getvalue(), scanner)
        from openpyxl import load_workbook
        wb2 = load_workbook(io.BytesIO(masked))
        assert "John Doe" not in str(wb2.worksheets[0]["A1"].value)
        assert "john.doe@org_name.com" not in str(wb2.worksheets[1]["A1"].value)


# -- PPTX tests ----------------------------------------------------------- #

class TestMaskPptx:
    def test_name_in_slide_masked(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx(["Presentation by John Doe"])
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        all_text = " ".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "John Doe" not in all_text
        assert "Presentation by" in all_text

    def test_no_pii_preserved(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx(["Q4 Business Review"])
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
        assert shapes[0].text_frame.text == "Q4 Business Review"

    def test_core_properties_masked(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx(["Content"], author="John Doe")
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        assert "John Doe" not in (prs.core_properties.author or "")

    def test_speaker_notes_masked(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx_with_notes(
            slide_text="Slide content",
            notes_text="Remind John Doe about the deadline",
        )
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        slide = prs.slides[0]
        assert slide.has_notes_slide
        notes = slide.notes_slide.notes_text_frame.text
        assert "John Doe" not in notes
        assert "deadline" in notes

    def test_table_in_slide_masked(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx_with_table([
            ["Team Member", "Role"],
            ["John Doe", "Lead"],
        ])
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        table = None
        for shape in prs.slides[0].shapes:
            if shape.has_table:
                table = shape.table
                break
        assert table is not None
        assert "John Doe" not in table.rows[1].cells[0].text
        assert table.rows[1].cells[1].text == "Lead"

    def test_multiple_slides(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        raw = _make_pptx([
            "Introduction",
            "Contact: john.doe@org_name.com",
            "Conclusion",
        ])
        masked = mask_pptx(raw, scanner)
        from pptx import Presentation
        prs = Presentation(io.BytesIO(masked))
        texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert texts[0] == "Introduction"
        assert "john.doe@org_name.com" not in texts[1]
        assert texts[2] == "Conclusion"


# -- is_office_doc --------------------------------------------------------- #

class TestIsOfficeDoc:
    def test_docx(self):
        from scripts.pii_mask.documents import is_office_doc
        assert is_office_doc("report.docx")
        assert is_office_doc("path/to/Report.DOCX")

    def test_xlsx(self):
        from scripts.pii_mask.documents import is_office_doc
        assert is_office_doc("data.xlsx")

    def test_pptx(self):
        from scripts.pii_mask.documents import is_office_doc
        assert is_office_doc("slides.pptx")

    def test_non_office(self):
        from scripts.pii_mask.documents import is_office_doc
        assert not is_office_doc("image.png")
        assert not is_office_doc("data.json")
        assert not is_office_doc("email.eml")
        assert not is_office_doc("report.pdf")
        assert not is_office_doc("legacy.doc")
        assert not is_office_doc("legacy.xls")


# -- Error handling -------------------------------------------------------- #

class TestErrorHandling:
    def test_corrupted_docx_raises(self, scanner):
        from scripts.pii_mask.documents import mask_docx
        with pytest.raises(Exception):
            mask_docx(b"not a docx file", scanner)

    def test_corrupted_xlsx_raises(self, scanner):
        from scripts.pii_mask.documents import mask_xlsx
        with pytest.raises(Exception):
            mask_xlsx(b"not an xlsx file", scanner)

    def test_corrupted_pptx_raises(self, scanner):
        from scripts.pii_mask.documents import mask_pptx
        with pytest.raises(Exception):
            mask_pptx(b"not a pptx file", scanner)
