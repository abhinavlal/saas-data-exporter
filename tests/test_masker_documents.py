"""Integration tests — Office document masking through per-service maskers."""

import io

import boto3
import pytest
from moto import mock_aws

from lib.s3 import S3Store
from scripts.pii_mask.pii_store import PIIStore
from scripts.pii_mask.scanner import TextScanner

SRC_BUCKET = "src-bucket"
DST_BUCKET = "dst-bucket"


# -- Fixtures -------------------------------------------------------------- #

@pytest.fixture
def store(tmp_path):
    s = PIIStore(str(tmp_path / "test.db"))
    s.add_domain("org_name.com", "example.com")
    s.get_or_create("EMAIL_ADDRESS", "john.doe@org_name.com")
    s.get_or_create("PERSON", "John Doe")
    return s


@pytest.fixture(scope="module")
def _analyzer():
    from presidio_analyzer import AnalyzerEngine
    return AnalyzerEngine()


@pytest.fixture
def scanner(store, _analyzer):
    s = TextScanner.__new__(TextScanner)
    s._store = store
    s._threshold = 0.5
    s._analyzer = _analyzer
    return s


@pytest.fixture
def s3_env():
    with mock_aws():
        conn = boto3.client("s3", region_name="us-east-1")
        conn.create_bucket(Bucket=SRC_BUCKET)
        conn.create_bucket(Bucket=DST_BUCKET)
        yield S3Store(bucket=SRC_BUCKET), S3Store(bucket=DST_BUCKET), conn


# -- Helpers --------------------------------------------------------------- #

def _make_docx(text: str) -> bytes:
    from docx import Document
    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(cell_value: str) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    wb.active["A1"] = cell_value
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1),
                                      Inches(5), Inches(1))
    txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# -- Google masker --------------------------------------------------------- #

class TestGoogleDocumentMasking:
    def test_drive_docx_masked(self, scanner, store, s3_env):
        from scripts.pii_mask.maskers.google import GoogleMasker

        src, dst, _ = s3_env
        masker = GoogleMasker(scanner)
        masked_email = store.lookup("EMAIL_ADDRESS", "john.doe@org_name.com")
        masked_slug = masked_email.replace("@", "_at_")

        docx_bytes = _make_docx("Report by John Doe")
        key = "google/john.doe_at_org_name.com/drive/abc123_Report.docx"
        src.upload_bytes(docx_bytes, key)

        status = masker.mask_file(src, dst, key)
        assert status == "ok"

        dst_key = f"google/{masked_slug}/drive/abc123_Report.docx"
        result = dst.download_bytes(dst_key)
        assert result is not None

        from docx import Document
        doc = Document(io.BytesIO(result))
        assert "John Doe" not in doc.paragraphs[0].text

    def test_gmail_attachment_xlsx_masked(self, scanner, store, s3_env):
        from scripts.pii_mask.maskers.google import GoogleMasker

        src, dst, _ = s3_env
        masker = GoogleMasker(scanner)
        masked_email = store.lookup("EMAIL_ADDRESS", "john.doe@org_name.com")
        masked_slug = masked_email.replace("@", "_at_")

        xlsx_bytes = _make_xlsx("john.doe@org_name.com")
        key = "google/john.doe_at_org_name.com/gmail/attachments/msg1/contacts.xlsx"
        src.upload_bytes(xlsx_bytes, key)

        status = masker.mask_file(src, dst, key)
        assert status == "ok"

        dst_key = f"google/{masked_slug}/gmail/attachments/msg1/contacts.xlsx"
        result = dst.download_bytes(dst_key)
        assert result is not None

        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(result))
        assert "john.doe@org_name.com" not in str(wb.active["A1"].value)

    def test_drive_keys_from_indexes(self, scanner, s3_env):
        """Drive index with downloaded office docs → keys include them."""
        from scripts.pii_mask.maskers.google import GoogleMasker

        src, dst, _ = s3_env
        masker = GoogleMasker(scanner, users={"john@org_name.com"})

        # Upload drive index with a docx entry
        drive_index = [
            {"id": "f1", "name": "Report.docx", "downloaded": True},
            {"id": "f2", "name": "photo.jpg", "downloaded": True},
            {"id": "f3", "name": "Data.xlsx", "downloaded": True},
            {"id": "f4", "name": "Skipped.docx", "downloaded": False},
        ]
        src.upload_json(drive_index,
                        "google/john_at_org_name.com/drive/_index.json")

        # Upload empty gmail/calendar indexes
        src.upload_json([], "google/john_at_org_name.com/gmail/_index.json")
        src.upload_json([], "google/john_at_org_name.com/calendar/_index.json")

        keys = masker.list_keys(src)

        # Should include office docs that were downloaded
        assert any("f1_Report.docx" in k for k in keys)
        assert any("f3_Data.xlsx" in k for k in keys)
        # Should NOT include non-office or not-downloaded
        assert not any("photo.jpg" in k for k in keys)
        assert not any("f4_Skipped.docx" in k for k in keys)

    def test_gmail_attachment_keys_from_indexes(self, scanner, s3_env):
        """Gmail index + S3 attachment listing → office doc keys."""
        from scripts.pii_mask.maskers.google import GoogleMasker

        src, dst, _ = s3_env
        masker = GoogleMasker(scanner, users={"john@org_name.com"})

        gmail_index = [{"id": "msg1"}, {"id": "msg2"}]
        src.upload_json(gmail_index,
                        "google/john_at_org_name.com/gmail/_index.json")
        src.upload_json([], "google/john_at_org_name.com/calendar/_index.json")
        src.upload_json(None, "google/john_at_org_name.com/drive/_index.json")

        # Upload some attachments
        src.upload_bytes(b"docx",
            "google/john_at_org_name.com/gmail/attachments/msg1/report.docx")
        src.upload_bytes(b"img",
            "google/john_at_org_name.com/gmail/attachments/msg1/logo.png")
        src.upload_bytes(b"pptx",
            "google/john_at_org_name.com/gmail/attachments/msg2/slides.pptx")

        keys = masker.list_keys(src)

        att_keys = [k for k in keys if "/attachments/" in k]
        assert any("report.docx" in k for k in att_keys)
        assert any("slides.pptx" in k for k in att_keys)
        assert not any("logo.png" in k for k in att_keys)


# -- Jira masker ---------------------------------------------------------- #

class TestJiraDocumentMasking:
    def test_attachment_docx_masked(self, scanner, s3_env):
        from scripts.pii_mask.maskers.jira import JiraMasker

        src, dst, _ = s3_env
        masker = JiraMasker(scanner)

        docx_bytes = _make_docx("Assigned to John Doe")
        key = "jira/PROJ/attachments/PROJ-1/spec.docx"
        src.upload_bytes(docx_bytes, key)

        status = masker.mask_file(src, dst, key)
        assert status == "ok"

        result = dst.download_bytes(key)
        assert result is not None

        from docx import Document
        doc = Document(io.BytesIO(result))
        assert "John Doe" not in doc.paragraphs[0].text

    def test_should_process_office_attachments(self, scanner):
        from scripts.pii_mask.maskers.jira import JiraMasker
        masker = JiraMasker(scanner)

        assert masker.should_process("jira/PROJ/attachments/PROJ-1/spec.docx")
        assert masker.should_process("jira/PROJ/attachments/PROJ-1/data.xlsx")
        assert masker.should_process("jira/PROJ/attachments/PROJ-1/deck.pptx")
        # Non-office attachments still skipped
        assert not masker.should_process("jira/PROJ/attachments/PROJ-1/image.png")
        assert not masker.should_process("jira/PROJ/attachments/PROJ-1/doc.pdf")
        # JSON tickets still processed
        assert masker.should_process("jira/PROJ/tickets/PROJ-1.json")

    def test_list_keys_includes_attachments(self, scanner, s3_env):
        from scripts.pii_mask.maskers.jira import JiraMasker

        src, dst, conn = s3_env
        masker = JiraMasker(scanner)

        # Upload ticket index and ticket JSON with attachment metadata
        ticket = {
            "key": "PROJ-1",
            "summary": "Test ticket",
            "attachments": [
                {"filename": "report.docx", "id": "att1"},
                {"filename": "screenshot.png", "id": "att2"},
                {"filename": "data.xlsx", "id": "att3"},
            ],
        }
        src.upload_json({"keys": ["PROJ-1"]}, "jira/PROJ/tickets/_index.json")
        src.upload_json(ticket, "jira/PROJ/tickets/PROJ-1.json")

        keys = masker.list_keys(src)

        assert any("report.docx" in k for k in keys)
        assert any("data.xlsx" in k for k in keys)
        assert not any("screenshot.png" in k for k in keys)


# -- Slack masker --------------------------------------------------------- #

class TestSlackDocumentMasking:
    def test_attachment_pptx_masked(self, scanner, s3_env):
        from scripts.pii_mask.maskers.slack import SlackMasker

        src, dst, _ = s3_env
        masker = SlackMasker(scanner)

        pptx_bytes = _make_pptx("Slides by John Doe")
        key = "slack/C123/attachments/F456_deck.pptx"
        src.upload_bytes(pptx_bytes, key)

        status = masker.mask_file(src, dst, key)
        assert status == "ok"

        result = dst.download_bytes(key)
        assert result is not None

        from pptx import Presentation
        prs = Presentation(io.BytesIO(result))
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                assert "John Doe" not in shape.text_frame.text

    def test_should_process_office_attachments(self, scanner):
        from scripts.pii_mask.maskers.slack import SlackMasker
        masker = SlackMasker(scanner)

        assert masker.should_process("slack/C123/attachments/F456_report.docx")
        assert masker.should_process("slack/C123/attachments/F456_data.xlsx")
        # Non-office attachments still skipped
        assert not masker.should_process("slack/C123/attachments/F456_image.png")
        # JSON messages still processed
        assert masker.should_process("slack/C123/messages/12345.json")

    def test_list_keys_includes_attachments(self, scanner, s3_env):
        from scripts.pii_mask.maskers.slack import SlackMasker

        src, dst, conn = s3_env
        masker = SlackMasker(scanner)

        # Upload channel structure
        src.upload_json({"name": "general"}, "slack/C123/channel_info.json")
        src.upload_json(["12345"], "slack/C123/messages/_index.json")
        src.upload_json({"text": "hello"}, "slack/C123/messages/12345.json")

        # Upload some attachments
        src.upload_bytes(b"docx", "slack/C123/attachments/F1_report.docx")
        src.upload_bytes(b"img", "slack/C123/attachments/F2_photo.png")

        keys = masker.list_keys(src)

        assert any("F1_report.docx" in k for k in keys)
        assert not any("F2_photo.png" in k for k in keys)


# -- Confluence masker ----------------------------------------------------- #

class TestConfluenceDocumentMasking:
    def test_should_process_office_attachments(self, scanner):
        from scripts.pii_mask.maskers.confluence import ConfluenceMasker
        masker = ConfluenceMasker(scanner)

        assert masker.should_process(
            "confluence/SPACE/attachments/123/report.docx")
        assert not masker.should_process(
            "confluence/SPACE/attachments/123/image.png")
        assert masker.should_process("confluence/SPACE/pages/123.json")


# -- BaseMasker ----------------------------------------------------------- #

class TestBaseMaskerDocumentSupport:
    def test_should_process_accepts_office_docs(self, scanner):
        from scripts.pii_mask.maskers.base import BaseMasker
        masker = BaseMasker(scanner)

        assert masker.should_process("any/path/file.docx")
        assert masker.should_process("any/path/file.xlsx")
        assert masker.should_process("any/path/file.pptx")
        assert masker.should_process("any/path/file.json")
        assert masker.should_process("any/path/file.eml")
        assert not masker.should_process("any/path/file.pdf")
        assert not masker.should_process("any/path/file.png")

    def test_mask_document_file_not_found(self, scanner, s3_env):
        from scripts.pii_mask.maskers.base import BaseMasker
        src, dst, _ = s3_env
        masker = BaseMasker(scanner)
        status = masker._mask_document_file(
            src, dst, "nonexistent/file.docx")
        assert status == "skipped (not found)"

    def test_mask_document_file_unsupported_ext(self, scanner, s3_env):
        from scripts.pii_mask.maskers.base import BaseMasker
        src, dst, _ = s3_env
        src.upload_bytes(b"data", "test/file.txt")
        masker = BaseMasker(scanner)
        status = masker._mask_document_file(src, dst, "test/file.txt")
        assert status == "skipped (unsupported ext)"

    def test_mask_document_file_corrupted(self, scanner, s3_env):
        from scripts.pii_mask.maskers.base import BaseMasker
        src, dst, _ = s3_env
        src.upload_bytes(b"not a real docx", "test/bad.docx")
        masker = BaseMasker(scanner)
        status = masker._mask_document_file(src, dst, "test/bad.docx")
        assert "error" in status
